"""
Módulo para cargar múltiples formatos de documentos en la plataforma AmershOp.
"""

import os
import json
import logging
from datetime import datetime
from typing import List

# Dependencias requeridas
from pypdf import PdfReader
import docx
import pandas as pd
from pptx import Presentation
from bs4 import BeautifulSoup
from langchain_core.documents import Document

# Configuración de registro (logging)
logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = [
    ".pdf", ".docx", ".xlsx", ".pptx", ".md", ".csv", ".json", ".html", ".htm"
]

def _get_base_metadata(file_path: str) -> dict:
    """Obtiene los metadatos base para un documento."""
    _, ext = os.path.splitext(file_path)
    return {
        "source": os.path.basename(file_path),
        "format": ext.lower()[1:] if ext else "",
        "file_path": file_path,
        "loaded_at": datetime.now().isoformat()
    }

def _load_pdf(file_path: str) -> List[Document]:
    """Carga el texto de un archivo PDF usando pypdf."""
    docs = []
    try:
        reader = PdfReader(file_path)
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                metadata = _get_base_metadata(file_path)
                metadata["page"] = i + 1
                docs.append(Document(page_content=text, metadata=metadata))
    except Exception as e:
        logger.error(f"Error al cargar PDF {file_path}: {e}")
    return docs

def _load_docx(file_path: str) -> List[Document]:
    """Carga el texto de un archivo DOCX."""
    docs = []
    try:
        doc = docx.Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        text = "\n".join(full_text)
        if text.strip():
            docs.append(Document(page_content=text, metadata=_get_base_metadata(file_path)))
    except Exception as e:
        logger.error(f"Error al cargar DOCX {file_path}: {e}")
    return docs

def _load_xlsx(file_path: str) -> List[Document]:
    """Carga datos de un archivo XLSX convirtiendo cada hoja en una tabla Markdown."""
    docs = []
    try:
        xls = pd.ExcelFile(file_path)
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            if not df.empty:
                text = df.to_markdown(index=False)
                metadata = _get_base_metadata(file_path)
                metadata["sheet"] = sheet_name
                docs.append(Document(page_content=text, metadata=metadata))
    except Exception as e:
        logger.error(f"Error al cargar XLSX {file_path}: {e}")
    return docs

def _load_pptx(file_path: str) -> List[Document]:
    """Carga el texto de un archivo PPTX extrayendo texto por diapositiva."""
    docs = []
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
            text = "\n".join(text_runs)
            if text.strip():
                metadata = _get_base_metadata(file_path)
                metadata["slide"] = i + 1
                docs.append(Document(page_content=text, metadata=metadata))
    except Exception as e:
        logger.error(f"Error al cargar PPTX {file_path}: {e}")
    return docs

def _load_markdown(file_path: str) -> List[Document]:
    """Carga el contenido de un archivo Markdown."""
    docs = []
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
        if text.strip():
            docs.append(Document(page_content=text, metadata=_get_base_metadata(file_path)))
    except UnicodeDecodeError:
        try:
            with open(file_path, "r", encoding="latin-1") as f:
                text = f.read()
            if text.strip():
                docs.append(Document(page_content=text, metadata=_get_base_metadata(file_path)))
        except Exception as e:
             logger.error(f"Error de codificación al cargar Markdown {file_path}: {e}")
    except Exception as e:
        logger.error(f"Error al cargar Markdown {file_path}: {e}")
    return docs

def _load_csv(file_path: str) -> List[Document]:
    """Carga un archivo CSV convirtiéndolo a una tabla Markdown."""
    docs = []
    try:
        df = pd.read_csv(file_path)
        if not df.empty:
            text = df.to_markdown(index=False)
            docs.append(Document(page_content=text, metadata=_get_base_metadata(file_path)))
    except Exception as e:
        logger.error(f"Error al cargar CSV {file_path}: {e}")
    return docs

def _load_json(file_path: str) -> List[Document]:
    """Carga un archivo JSON aplanando los datos en texto legible."""
    docs = []
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        text = json.dumps(data, indent=2, ensure_ascii=False)
        if text.strip():
            docs.append(Document(page_content=text, metadata=_get_base_metadata(file_path)))
    except Exception as e:
        logger.error(f"Error al cargar JSON {file_path}: {e}")
    return docs

def _load_html(file_path: str) -> List[Document]:
    """Carga un archivo HTML extrayendo el texto limpio."""
    docs = []
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f, "html.parser")
            text = soup.get_text(separator="\n", strip=True)
        if text.strip():
            docs.append(Document(page_content=text, metadata=_get_base_metadata(file_path)))
    except Exception as e:
        logger.error(f"Error al cargar HTML {file_path}: {e}")
    return docs

def load_single_document(file_path: str) -> List[Document]:
    """
    Detecta el formato del archivo por su extensión y carga su contenido correspondiente.
    
    Args:
        file_path (str): Ruta al archivo a cargar.
        
    Returns:
        List[Document]: Lista de documentos cargados.
    """
    if not os.path.exists(file_path):
        logger.error(f"El archivo no existe: {file_path}")
        return []
        
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    
    if ext == ".pdf":
        return _load_pdf(file_path)
    elif ext == ".docx":
        return _load_docx(file_path)
    elif ext == ".xlsx":
        return _load_xlsx(file_path)
    elif ext == ".pptx":
        return _load_pptx(file_path)
    elif ext == ".md":
        return _load_markdown(file_path)
    elif ext == ".csv":
        return _load_csv(file_path)
    elif ext == ".json":
        return _load_json(file_path)
    elif ext in [".html", ".htm"]:
        return _load_html(file_path)
    else:
        logger.warning(f"Extensión no soportada para el archivo: {file_path}")
        return []

def load_all_documents(directory: str) -> List[Document]:
    """
    Carga todos los archivos soportados desde un directorio.
    
    Args:
        directory (str): Ruta al directorio.
        
    Returns:
        List[Document]: Lista de documentos cargados.
    """
    all_docs = []
    if not os.path.isdir(directory):
        logger.error(f"El directorio no existe: {directory}")
        return all_docs
        
    for root, _, files in os.walk(directory):
        for file in files:
            _, ext = os.path.splitext(file)
            if ext.lower() in SUPPORTED_EXTENSIONS:
                file_path = os.path.join(root, file)
                all_docs.extend(load_single_document(file_path))
                
    return all_docs
