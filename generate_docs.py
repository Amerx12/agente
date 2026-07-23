import os
import json
import csv
from datetime import datetime

# Dependencias externas necesarias:
# pip install reportlab python-docx openpyxl python-pptx

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.enums import TA_JUSTIFY
except ImportError:
    print("Por favor instala reportlab: pip install reportlab")

try:
    from docx import Document
    from docx.shared import Pt, RGBColor
except ImportError:
    print("Por favor instala python-docx: pip install python-docx")

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill
except ImportError:
    print("Por favor instala openpyxl: pip install openpyxl")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Por favor instala python-pptx: pip install python-pptx")

BASE_DIR = r"c:\Users\amer\Desktop\agente\documents"

def crear_directorio_base():
    """Crea el directorio base si no existe."""
    if not os.path.exists(BASE_DIR):
        os.makedirs(BASE_DIR)
        print(f"Directorio creado: {BASE_DIR}")

def generar_politica_privacidad():
    """Genera politica_privacidad.pdf usando reportlab."""
    filename = os.path.join(BASE_DIR, "politica_privacidad.pdf")
    doc = SimpleDocTemplate(filename, pagesize=letter)
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name='Justify', alignment=TA_JUSTIFY, parent=styles['Normal']))
    
    Story = []
    
    Story.append(Paragraph("Política de Privacidad de AmershOp", styles['Title']))
    Story.append(Spacer(1, 12))
    
    # Add substantial content here
    content = [
        ("1. Introducción", "En AmershOp, valoramos su privacidad y nos comprometemos a proteger sus datos personales. Esta Política de Privacidad explica cómo recopilamos, usamos, compartimos y protegemos su información cuando visita nuestro sitio web o realiza compras."),
        ("2. Datos Recopilados", "Recopilamos diferentes tipos de información, incluyendo:\n- Datos personales: Nombre, dirección, correo electrónico, número de teléfono.\n- Datos de navegación: Dirección IP, tipo de navegador, páginas visitadas.\n- Datos de pago: Información de tarjeta de crédito/débito, historial de compras."),
        ("3. Uso de Cookies", "Utilizamos cookies y tecnologías similares para mejorar su experiencia de navegación, personalizar el contenido y los anuncios, y analizar nuestro tráfico. Puede configurar su navegador para rechazar las cookies, pero esto puede limitar algunas funciones de nuestro sitio."),
        ("4. Derechos ARCO", "Usted tiene derecho a Acceder, Rectificar, Cancelar u Oponerse al tratamiento de sus datos personales. Para ejercer estos derechos, comuníquese con nosotros a través de los canales proporcionados al final de este documento."),
        ("5. Compartición con Terceros", "No vendemos ni alquilamos sus datos a terceros. Podemos compartir su información con proveedores de servicios de confianza que nos ayudan a operar nuestro negocio, procesar pagos o entregar pedidos, siempre bajo estrictos acuerdos de confidencialidad."),
        ("6. Retención de Datos", "Conservaremos sus datos personales durante el tiempo que sea necesario para cumplir con los propósitos descritos en esta política, a menos que la ley exija o permita un período de retención más largo."),
        ("7. Seguridad de los Datos", "Implementamos medidas de seguridad técnicas y organizativas para proteger sus datos contra acceso no autorizado, pérdida o alteración. Sin embargo, ninguna transmisión por Internet es 100% segura."),
        ("8. Contacto del DPO", "Si tiene preguntas sobre nuestra Política de Privacidad o el tratamiento de sus datos, comuníquese con nuestro Oficial de Protección de Datos (DPO) en dpo@amershop.com o llame al 555-0198.")
    ]
    
    # Adding text multiple times to ensure length for testing
    for title, text in content:
        Story.append(Paragraph(title, styles['Heading2']))
        Story.append(Spacer(1, 6))
        Story.append(Paragraph(text, styles['Justify']))
        Story.append(Spacer(1, 12))
        
        # Add filler text to make it realistic 2 pages
        filler = "Adicionalmente, " + text.lower() + " Esto asegura el cumplimiento de las normativas internacionales de protección de datos como el GDPR y leyes locales aplicables, manteniendo siempre el más alto estándar de seguridad y transparencia para con nuestros valiosos clientes. Nuestro compromiso es continuo y sujeto a revisiones regulares."
        Story.append(Paragraph(filler, styles['Justify']))
        Story.append(Spacer(1, 12))

    doc.build(Story)
    print(f"Generado: {filename}")

def generar_politica_reembolso():
    """Genera politica_reembolso.docx usando python-docx."""
    filename = os.path.join(BASE_DIR, "politica_reembolso.docx")
    doc = Document()
    
    doc.add_heading('Política de Reembolsos y Devoluciones - AmershOp', 0)
    
    doc.add_paragraph('En AmershOp queremos que estés 100% satisfecho con tu compra. Si no estás conforme, te explicamos cómo proceder.')
    
    doc.add_heading('Plazos de Devolución', level=1)
    doc.add_paragraph('Dispones de 30 días naturales desde la recepción de tu pedido para solicitar una devolución o reembolso completo del producto adquirido en nuestra plataforma. Tras este periodo, lamentablemente no podremos ofrecerte un reembolso o cambio.')
    
    doc.add_heading('Condiciones del Producto', level=1)
    p = doc.add_paragraph('Para ser elegible para una devolución, el artículo debe:')
    p.add_run('\n1. Estar sin usar y en las mismas condiciones en que lo recibiste.').bold = True
    p.add_run('\n2. Conservar el embalaje original sin daños severos.').bold = True
    p.add_run('\n3. Incluir todos los manuales, accesorios y regalos promocionales.').bold = True
    
    doc.add_heading('Proceso Paso a Paso', level=1)
    doc.add_paragraph('1. Inicia sesión en tu cuenta y ve a "Mis Pedidos".\n2. Selecciona el pedido y el producto a devolver.\n3. Elige el motivo de la devolución y si prefieres reembolso o cambio.\n4. Imprime la etiqueta de envío prepagada que enviaremos a tu correo.\n5. Empaqueta el producto y entrégalo en cualquier sucursal de nuestros transportistas asociados.')
    
    doc.add_heading('Excepciones', level=1)
    doc.add_paragraph('Los siguientes artículos no son retornables:\n- Software descargable o tarjetas de regalo.\n- Productos de higiene íntima o cuidado personal.\n- Artículos en liquidación final o con descuento mayor al 50% (a menos que presenten defecto de fábrica).')
    
    doc.add_heading('Reembolso vs Cambio', level=1)
    doc.add_paragraph('Una vez recibida e inspeccionada tu devolución, te enviaremos un correo electrónico notificándote la aprobación o rechazo. Si es aprobado, el reembolso se procesará y se aplicará automáticamente a tu tarjeta de crédito o método de pago original en un plazo de 5 a 10 días hábiles. Si elegiste cambio, el nuevo producto será despachado inmediatamente.')
    
    doc.add_heading('Gastos de Envío en Devolución', level=1)
    doc.add_paragraph('Si la devolución es por defecto de fábrica o error nuestro, AmershOp cubrirá el 100% de los gastos de envío. Si la devolución es por cambio de opinión o error del cliente al elegir el producto, se descontará una tarifa plana de $150 MXN del total a reembolsar por concepto de envío.')
    
    doc.save(filename)
    print(f"Generado: {filename}")

def generar_faq():
    """Genera faq.md."""
    filename = os.path.join(BASE_DIR, "faq.md")
    content = """# Preguntas Frecuentes (FAQ) - AmershOp

## Cuenta y Accesos
**1. ¿Cómo creo una cuenta en AmershOp?**
Para crear una cuenta, haz clic en el botón "Registrarse" en la esquina superior derecha, ingresa tu correo electrónico, crea una contraseña segura y completa tus datos personales.

**2. ¿Qué hago si olvidé mi contraseña?**
Ve a "Iniciar Sesión" y selecciona "Olvidé mi contraseña". Ingresa tu correo y te enviaremos un enlace seguro para restablecerla.

**3. ¿Puedo comprar como invitado sin crear cuenta?**
Sí, puedes realizar compras como invitado, pero crear una cuenta te permite rastrear tus pedidos fácilmente y guardar métodos de pago.

## Métodos de Pago
**4. ¿Qué métodos de pago aceptan?**
Aceptamos tarjetas de crédito y débito (Visa, MasterCard, American Express), PayPal, transferencias bancarias SPEI, y pagos en efectivo en tiendas de conveniencia (OXXO).

**5. ¿Es seguro ingresar mi tarjeta de crédito?**
Absolutamente. Utilizamos encriptación SSL de 256 bits y nuestros procesos de pago cumplen con los estándares PCI-DSS. No almacenamos los datos completos de tu tarjeta.

**6. ¿Puedo pagar a meses sin intereses (MSI)?**
Sí, ofrecemos hasta 12 MSI en compras superiores a $2,500 MXN con tarjetas de crédito participantes.

## Tiempos de Envío
**7. ¿Cuánto tiempo tarda en llegar mi pedido?**
El envío estándar toma de 3 a 5 días hábiles. El envío express toma de 1 a 2 días hábiles dependiendo de tu código postal.

**8. ¿Desde dónde se envían los productos?**
Contamos con múltiples centros de distribución a nivel nacional. Tu pedido se enviará desde el almacén más cercano con disponibilidad de stock.

**9. ¿Realizan envíos internacionales?**
Actualmente, solo realizamos envíos dentro del territorio nacional. Estamos trabajando para expandir nuestras fronteras pronto.

## Garantías
**10. ¿Todos los productos tienen garantía?**
Sí, todos nuestros productos cuentan con una garantía mínima de 1 año contra defectos de fabricación directamente con el proveedor, respaldada por AmershOp durante los primeros 30 días.

**11. ¿Cómo hago válida una garantía?**
Contacta a nuestro equipo de soporte técnico con tu número de pedido y evidencia del fallo (fotos o video). Te guiaremos en el proceso.

**12. ¿Qué invalida la garantía?**
La garantía no cubre daños por mal uso, golpes, derrames de líquidos, modificaciones no autorizadas o desgaste natural.

## Cambios y Devoluciones
**13. ¿Tengo tiempo límite para devolver un producto?**
Sí, tienes 30 días naturales a partir de la recepción del pedido para solicitar una devolución.

**14. ¿Las devoluciones tienen algún costo?**
Si el producto es defectuoso, la devolución es gratuita. Si es por cambio de opinión, se deducirá una tarifa de envío de $150 MXN.

**15. ¿Cuánto tarda el reembolso?**
Una vez aprobado, el reembolso tarda entre 5 y 10 días hábiles en reflejarse en tu estado de cuenta, dependiendo de tu banco.

## Facturación
**16. ¿Cómo solicito la factura de mi compra?**
Durante el checkout, marca la casilla "Requiero Factura" e ingresa tus datos fiscales. También puedes solicitarla en la sección "Mis Pedidos" hasta fin de mes.

**17. ¿Cuándo recibiré mi factura?**
La factura en formato PDF y XML será enviada a tu correo electrónico en un plazo máximo de 24 horas hábiles tras la confirmación del pago.

**18. Me equivoqué en mis datos fiscales, ¿puedo refacturar?**
Sí, tienes 72 horas para solicitar la corrección de una factura contactando a nuestro equipo de atención al cliente.

## Soporte Técnico
**19. ¿Ofrecen ayuda para instalar los productos?**
Brindamos soporte técnico remoto básico. Para instalación física de equipos complejos, te recomendamos consultar con técnicos locales certificados.

**20. ¿Cómo contacto al equipo de soporte?**
Puedes contactarnos vía chat en vivo de lunes a viernes (9am-6pm), por correo electrónico a soporte@amershop.com, o llamando al 800-AMERSHOP.
"""
    with open(filename, 'w', encoding='utf-8') as f:
        f.write(content)
    print(f"Generado: {filename}")

def generar_guia_envios():
    """Genera guia_envios.html."""
    filename = os.path.join(BASE_DIR, "guia_envios.html")
    content = """<!DOCTYPE html>
<html lang="es">
<head>
    <meta charset="UTF-8">
    <title>Guía de Envíos - AmershOp</title>
    <style>
        body { font-family: Arial, sans-serif; line-height: 1.6; margin: 40px; color: #333; }
        h1, h2 { color: #0056b3; }
        table { width: 100%; border-collapse: collapse; margin-bottom: 20px; }
        th, td { border: 1px solid #ddd; padding: 12px; text-align: left; }
        th { background-color: #f2f2f2; color: #333; }
        .highlight { background-color: #e6f7ff; padding: 15px; border-left: 5px solid #0056b3; }
    </style>
</head>
<body>
    <h1>Guía Completa de Envíos AmershOp</h1>
    <p>En AmershOp nos esforzamos por entregar tus productos tecnológicos de la manera más rápida y segura posible. A continuación, detallamos toda la información que necesitas saber sobre nuestras políticas de envío.</p>
    
    <h2>1. Zonas de Cobertura</h2>
    <p>Nuestra red logística abarca las siguientes zonas:</p>
    <ul>
        <li><strong>Local:</strong> Área metropolitana y municipios conurbados de nuestra sede principal.</li>
        <li><strong>Nacional:</strong> Resto de los estados de la república.</li>
        <li><strong>Internacional:</strong> (Próximamente) Norteamérica y Centroamérica.</li>
    </ul>

    <h2>2. Tiempos de Entrega y Costos</h2>
    <p>Los tiempos de entrega son estimados y comienzan a contar a partir de la confirmación del pago.</p>
    <table>
        <tr>
            <th>Tipo de Envío</th>
            <th>Tiempo Estimado</th>
            <th>Costo Base</th>
            <th>Descripción</th>
        </tr>
        <tr>
            <td>Envío Estándar Local</td>
            <td>1 - 2 días hábiles</td>
            <td>$89 MXN</td>
            <td>Entrega económica en el área local.</td>
        </tr>
        <tr>
            <td>Envío Estándar Nacional</td>
            <td>3 - 5 días hábiles</td>
            <td>$129 MXN</td>
            <td>Entrega confiable a nivel nacional vía paqueterías socias.</td>
        </tr>
        <tr>
            <td>Envío Express Local</td>
            <td>Mismo día (pedidos antes de 12pm)</td>
            <td>$149 MXN</td>
            <td>Entrega urgente en motocicleta o van dedicada.</td>
        </tr>
        <tr>
            <td>Envío Express Nacional</td>
            <td>1 - 2 días hábiles</td>
            <td>$249 MXN</td>
            <td>Envío aéreo prioritario a nivel nacional.</td>
        </tr>
    </table>

    <div class="highlight">
        <h3>¡Envío Gratuito!</h3>
        <p>Disfruta de <strong>ENVÍO ESTÁNDAR GRATIS</strong> en todas tus compras superiores a <strong>$999 MXN</strong>. El descuento se aplicará automáticamente en el carrito de compras.</p>
    </div>

    <h2>3. Rastreo de Pedidos</h2>
    <p>Una vez que tu pedido sea despachado, recibirás un correo electrónico con un número de guía de 10 dígitos y un enlace al portal del transportista para rastrear el paquete en tiempo real.</p>
    
    <h2>4. Consideraciones Importantes</h2>
    <ul>
        <li>Los pedidos realizados en fines de semana o días festivos se procesarán el siguiente día hábil.</li>
        <li>Si no estás en casa durante el intento de entrega, el transportista dejará un aviso e intentará nuevamente al día siguiente. (Máximo 3 intentos).</li>
        <li>Para zonas de difícil acceso ("zonas extendidas"), el tiempo de entrega puede incrementar de 2 a 3 días adicionales.</li>
    </ul>
</body>
</html>
"""
    with open(filename, 'w', encoding='utf-8') as f:
        f.write(content)
    print(f"Generado: {filename}")

def generar_terminos_condiciones():
    """Genera terminos_condiciones.pdf usando reportlab."""
    filename = os.path.join(BASE_DIR, "terminos_condiciones.pdf")
    doc = SimpleDocTemplate(filename, pagesize=letter)
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name='Justify', alignment=TA_JUSTIFY, parent=styles['Normal']))
    
    Story = []
    Story.append(Paragraph("Términos y Condiciones de Uso - AmershOp", styles['Title']))
    Story.append(Spacer(1, 12))
    
    content = [
        ("1. Uso del Sitio Web", "Bienvenido a AmershOp. Al acceder y utilizar este sitio web, usted acepta cumplir y estar sujeto a estos Términos y Condiciones. Si no está de acuerdo con alguna parte, no debe utilizar nuestro sitio. El contenido es para su uso personal y no comercial."),
        ("2. Registro de Cuenta", "Para acceder a ciertas funciones, puede requerirse crear una cuenta. Usted es responsable de mantener la confidencialidad de su contraseña y asume la responsabilidad de todas las actividades bajo su cuenta. Nos reservamos el derecho de cancelar cuentas sospechosas de fraude."),
        ("3. Precios y Disponibilidad", "Todos los precios están sujetos a cambios sin previo aviso. Nos esforzamos por mantener el inventario actualizado; sin embargo, en caso de que un producto no esté disponible tras realizar el pedido, le notificaremos de inmediato y procederemos con el reembolso."),
        ("4. Proceso de Compra", "Al realizar un pedido, usted hace una oferta de compra. Recibirá un correo de confirmación, lo cual no garantiza la aceptación del pedido. Nos reservamos el derecho de cancelar o limitar cantidades por cliente, cuenta o tarjeta de crédito."),
        ("5. Garantías de Productos", "La mayoría de nuestros artículos incluyen una garantía del fabricante. AmershOp actúa como facilitador para procesar garantías durante los primeros 30 días; posteriormente, deberá contactar directamente al fabricante según sus propias políticas."),
        ("6. Limitación de Responsabilidad", "AmershOp no será responsable por daños indirectos, incidentales, especiales o consecuentes que surjan del uso o la imposibilidad de usar nuestros productos o el sitio web. Nuestra responsabilidad máxima se limitará al precio pagado por el producto."),
        ("7. Propiedad Intelectual", "Todo el contenido del sitio, incluyendo textos, gráficos, logos e imágenes, es propiedad de AmershOp o de sus proveedores, y está protegido por leyes de derechos de autor internacionales."),
        ("8. Ley Aplicable", "Estos términos se regirán e interpretarán de acuerdo con las leyes del país en que operamos comercialmente, sin dar efecto a ningún principio de conflictos de leyes. Cualquier disputa se someterá a la jurisdicción exclusiva de los tribunales competentes en nuestra sede principal.")
    ]
    
    for title, text in content:
        Story.append(Paragraph(title, styles['Heading2']))
        Story.append(Spacer(1, 6))
        Story.append(Paragraph(text, styles['Justify']))
        Story.append(Spacer(1, 12))
        
        # Filler text to extend to 2 pages
        filler = "Adicionalmente, " + text.lower() + " Esto constituye un acuerdo vinculante entre las partes. La empresa se reserva el derecho de modificar estos términos en cualquier momento. El uso continuado del sitio tras dichas modificaciones constituirá su reconocimiento y aceptación de los nuevos términos. Le instamos a revisar esta sección periódicamente para mantenerse informado."
        Story.append(Paragraph(filler, styles['Justify']))
        Story.append(Spacer(1, 12))

    doc.build(Story)
    print(f"Generado: {filename}")

def generar_catalogo_productos():
    """Genera catalogo_productos.xlsx usando openpyxl."""
    filename = os.path.join(BASE_DIR, "catalogo_productos.xlsx")
    wb = Workbook()
    ws = wb.active
    ws.title = "Catálogo"
    
    headers = ["ID", "Nombre", "Categoría", "Marca", "Precio", "Stock", "Descripción"]
    ws.append(headers)
    
    # Styling headers
    header_fill = PatternFill(start_color="0056b3", end_color="0056b3", fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True)
    for col in range(1, 8):
        cell = ws.cell(row=1, column=col)
        cell.fill = header_fill
        cell.font = header_font
    
    # Generate 30 products
    products = [
        ("LP-001", "Laptop Pro X15", "Laptops", "TechBrand", 25999, 15, "Laptop 15.6'', Intel Core i7, 16GB RAM, 512GB SSD"),
        ("LP-002", "Laptop Ultra Lite", "Laptops", "UltraSys", 18500, 22, "Laptop ultraligera 13'', AMD Ryzen 5, 8GB RAM"),
        ("LP-003", "Gaming Beast V2", "Laptops", "GamerX", 32000, 8, "Laptop gaming RTX 3070, 32GB RAM, Pantalla 144Hz"),
        ("LP-004", "Office Book A1", "Laptops", "TechBrand", 12000, 45, "Ideal para oficina y estudiantes. Celeron, 4GB RAM"),
        ("LP-005", "Creator Studio P", "Laptops", "VisionTech", 45000, 5, "Pantalla OLED 4K, ideal para diseñadores. RTX 4060"),
        ("SM-001", "Phone Alpha 12", "Smartphones", "MobileCo", 15999, 50, "Smartphone 5G, 128GB, cámara triple de 64MP"),
        ("SM-002", "Phone Alpha 12 Pro", "Smartphones", "MobileCo", 19999, 30, "Versión Pro, 256GB, Zoom óptico 5x"),
        ("SM-003", "Note Master 5", "Smartphones", "StylusM", 22500, 15, "Incluye lápiz óptico, pantalla AMOLED 6.8''"),
        ("SM-004", "Eco Phone Basic", "Smartphones", "EcoTech", 4500, 100, "Batería de larga duración, 32GB, Android Go"),
        ("SM-005", "Foldable Z1", "Smartphones", "FutureM", 35000, 10, "Pantalla plegable, tecnología de vanguardia"),
        ("AC-001", "Mouse Inalámbrico", "Accesorios", "ClickIt", 450, 200, "Mouse ergonómico 2.4GHz, batería incluida"),
        ("AC-002", "Teclado Mecánico RGB", "Accesorios", "TypePro", 1200, 80, "Switches azules, retroiluminación personalizable"),
        ("AC-003", "Hub USB-C 7 en 1", "Accesorios", "ConnectX", 850, 120, "HDMI, USB 3.0, lector SD/TF, carga PD"),
        ("AC-004", "Funda Laptop 15''", "Accesorios", "SafeBags", 350, 150, "Neopreno resistente al agua, color negro"),
        ("AC-005", "Soporte Ajustable Aluminio", "Accesorios", "StandIt", 650, 60, "Soporte para laptop o tablet, mejora postura"),
        ("AU-001", "Audífonos Inalámbricos Pro", "Audio", "SoundMax", 2500, 40, "Cancelación activa de ruido, 24h de batería"),
        ("AU-002", "Bocina Bluetooth", "Audio", "BassBoom", 1200, 75, "Resistente al agua IPX7, 20W de potencia"),
        ("AU-003", "Earbuds Sport", "Audio", "FitSound", 899, 100, "Diseño deportivo, agarre seguro, sudor-resistente"),
        ("AU-004", "Micrófono USB Podcast", "Audio", "VocalPro", 1800, 25, "Micrófono condensador cardioide con soporte"),
        ("AU-005", "Barra de Sonido TV", "Audio", "HomeCinema", 3500, 12, "Sistema 2.1 con subwoofer inalámbrico"),
        ("GM-001", "Control Inalámbrico", "Gaming", "PlayTech", 1100, 90, "Compatible con PC y consolas, vibración"),
        ("GM-002", "Silla Gamer Ergonómica", "Gaming", "SitPro", 4500, 20, "Reclinable 180 grados, cojines lumbar y cervical"),
        ("GM-003", "Monitor 27'' 165Hz", "Gaming", "VisionTech", 6500, 15, "Panel IPS, 1ms respuesta, FreeSync"),
        ("GM-004", "Diadema Gamer 7.1", "Gaming", "SoundMax", 1400, 50, "Sonido envolvente virtual, micrófono retráctil"),
        ("GM-005", "Mousepad RGB XXL", "Gaming", "DeskMats", 500, 85, "Tamaño 90x40cm, bordes iluminados"),
        ("TB-001", "Tab Pro 11", "Tablets", "TechBrand", 8500, 35, "Tablet 11 pulgadas, 128GB, Octa-core"),
        ("TB-002", "Tab Mini 8", "Tablets", "TechBrand", 4500, 60, "Compacta, ideal para lectura y multimedia"),
        ("TB-003", "Pad Creator 12.9", "Tablets", "VisionTech", 18000, 18, "Orientada a diseño gráfico, incluye stylus"),
        ("TB-004", "Kids Pad Segura", "Tablets", "EduTech", 2500, 45, "Funda protectora de goma, control parental"),
        ("TB-005", "Tab Work E-ink", "Tablets", "ReadPro", 9500, 10, "Pantalla de tinta electrónica, para notas y lectura")
    ]
    
    for p in products:
        ws.append(p)
        
    wb.save(filename)
    print(f"Generado: {filename}")

def generar_precios_envio():
    """Genera precios_envio.csv."""
    filename = os.path.join(BASE_DIR, "precios_envio.csv")
    headers = ["Zona", "Peso_Min_kg", "Peso_Max_kg", "Costo_Estandar", "Costo_Express", "Tiempo_Estandar_Dias", "Tiempo_Express_Dias"]
    
    data = [
        ["Local", 0.0, 1.0, 89.0, 149.0, 2, 1],
        ["Local", 1.01, 5.0, 119.0, 189.0, 2, 1],
        ["Local", 5.01, 15.0, 159.0, 259.0, 2, 1],
        ["Local", 15.01, 30.0, 250.0, 400.0, 3, 2],
        ["Local", 30.01, 999.0, 400.0, 650.0, 3, 2],
        
        ["Nacional_Cercano", 0.0, 1.0, 129.0, 249.0, 3, 1],
        ["Nacional_Cercano", 1.01, 5.0, 169.0, 329.0, 3, 2],
        ["Nacional_Cercano", 5.01, 15.0, 229.0, 459.0, 4, 2],
        ["Nacional_Cercano", 15.01, 30.0, 350.0, 700.0, 5, 3],
        ["Nacional_Cercano", 30.01, 999.0, 600.0, 1200.0, 5, 4],
        
        ["Nacional_Lejano", 0.0, 1.0, 159.0, 299.0, 5, 2],
        ["Nacional_Lejano", 1.01, 5.0, 209.0, 399.0, 5, 2],
        ["Nacional_Lejano", 5.01, 15.0, 289.0, 559.0, 6, 3],
        ["Nacional_Lejano", 15.01, 30.0, 450.0, 900.0, 7, 4],
        ["Nacional_Lejano", 30.01, 999.0, 800.0, 1500.0, 8, 5],
    ]
    
    with open(filename, 'w', newline='', encoding='utf-8') as f:
        writer = csv.writer(f)
        writer.writerow(headers)
        writer.writerows(data)
    print(f"Generado: {filename}")

def generar_config_tienda():
    """Genera config_tienda.json."""
    filename = os.path.join(BASE_DIR, "config_tienda.json")
    config = {
        "tienda": {
            "nombre": "AmershOp",
            "url": "https://www.amershop.com",
            "email_contacto": "contacto@amershop.com",
            "telefono": "+52-800-AMERSHOP",
            "horario_atencion": "Lunes a Viernes 09:00 - 18:00, Sábados 10:00 - 14:00"
        },
        "configuracion_regional": {
            "moneda": "MXN",
            "simbolo_moneda": "$",
            "impuesto_iva_porcentaje": 16,
            "idioma_por_defecto": "es-MX"
        },
        "pagos_aceptados": [
            "Tarjeta de Crédito",
            "Tarjeta de Débito",
            "PayPal",
            "Transferencia SPEI",
            "Pago en OXXO"
        ],
        "redes_sociales": {
            "facebook": "https://facebook.com/amershop_mx",
            "instagram": "https://instagram.com/amershop_tech",
            "twitter": "https://twitter.com/amershop_oficial"
        },
        "configuracion_envios": {
            "umbral_envio_gratis": 999.00,
            "transportistas_activos": ["FedEx", "DHL", "Estafeta", "Mensajería Local"]
        },
        "politicas": {
            "dias_limite_devolucion": 30,
            "dias_garantia_base": 365,
            "costo_reestocaje_devolucion": 150.00
        }
    }
    with open(filename, 'w', encoding='utf-8') as f:
        json.dump(config, f, indent=4, ensure_ascii=False)
    print(f"Generado: {filename}")

def generar_presentacion_empresa():
    """Genera presentacion_empresa.pptx usando python-pptx."""
    filename = os.path.join(BASE_DIR, "presentacion_empresa.pptx")
    prs = Presentation()
    
    # Slide 1: Título
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AmershOp - Tecnología para Todos"
    subtitle.text = "Presentación Corporativa\\n" + datetime.now().strftime("%Y")
    
    # Slide 2: Misión y Visión
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide2.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Misión y Visión"
    tf = body_shape.text_frame
    tf.text = "Misión:"
    p = tf.add_paragraph()
    p.text = "Democratizar el acceso a la tecnología de vanguardia ofreciendo productos de alta calidad a precios competitivos con un servicio excepcional."
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "Visión:"
    p3 = tf.add_paragraph()
    p3.text = "Ser la plataforma de comercio electrónico líder en productos tecnológicos en América Latina, reconocida por la confianza y satisfacción de sus usuarios."
    p3.level = 1
    
    # Slide 3: Valores Corporativos
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    shapes3 = slide3.shapes
    title_shape3 = shapes3.title
    body_shape3 = shapes3.placeholders[1]
    title_shape3.text = "Valores Corporativos"
    tf3 = body_shape3.text_frame
    valores = ["Innovación constante", "Transparencia total", "Pasión por el cliente", "Calidad garantizada", "Sostenibilidad"]
    for val in valores:
        p = tf3.add_paragraph()
        p.text = val
        p.level = 0
        
    # Slide 4: Nuestro equipo (departamentos)
    slide4 = prs.slides.add_slide(bullet_slide_layout)
    shapes4 = slide4.shapes
    title_shape4 = shapes4.title
    body_shape4 = shapes4.placeholders[1]
    title_shape4.text = "Nuestro Equipo"
    tf4 = body_shape4.text_frame
    deps = [
        "Tecnología (Desarrollo y Sistemas)", 
        "Operaciones y Logística", 
        "Atención al Cliente (Soporte 24/7)", 
        "Marketing y Ventas", 
        "Compras y Selección de Producto"
    ]
    for dep in deps:
        p = tf4.add_paragraph()
        p.text = dep
        p.level = 0
        
    # Slide 5: Cifras clave
    slide5 = prs.slides.add_slide(bullet_slide_layout)
    shapes5 = slide5.shapes
    title_shape5 = shapes5.title
    body_shape5 = shapes5.placeholders[1]
    title_shape5.text = "Cifras Clave del Último Año"
    tf5 = body_shape5.text_frame
    cifras = [
        "+500,000 pedidos entregados",
        "+1.2 millones de usuarios activos mensuales",
        "Catálogo de más de 10,000 productos",
        "Tasa de satisfacción del cliente: 98%",
        "Presencia en toda la república"
    ]
    for c in cifras:
        p = tf5.add_paragraph()
        p.text = c
        p.level = 0
        
    # Slide 6: Canales de contacto
    slide6 = prs.slides.add_slide(bullet_slide_layout)
    shapes6 = slide6.shapes
    title_shape6 = shapes6.title
    body_shape6 = shapes6.placeholders[1]
    title_shape6.text = "Canales de Contacto"
    tf6 = body_shape6.text_frame
    contactos = [
        "Sitio Web: www.amershop.com",
        "Correo: corporativo@amershop.com",
        "Teléfono: 800-AMERSHOP",
        "Oficinas: Av. Tecnológica 100, Piso 5, Ciudad Central",
        "Redes Sociales: @amershop_oficial"
    ]
    for c in contactos:
        p = tf6.add_paragraph()
        p.text = c
        p.level = 0
        
    prs.save(filename)
    print(f"Generado: {filename}")

def generate_all():
    """Genera todos los documentos de muestra."""
    print("Iniciando generación de documentos AmershOp...")
    crear_directorio_base()
    
    # Manejamos errores en caso de que falten librerías para no romper toda la ejecución
    try:
        generar_politica_privacidad()
    except Exception as e:
        print(f"Error generando politica_privacidad.pdf: {e}")
        
    try:
        generar_politica_reembolso()
    except Exception as e:
        print(f"Error generando politica_reembolso.docx: {e}")
        
    try:
        generar_faq()
    except Exception as e:
        print(f"Error generando faq.md: {e}")
        
    try:
        generar_guia_envios()
    except Exception as e:
        print(f"Error generando guia_envios.html: {e}")
        
    try:
        generar_terminos_condiciones()
    except Exception as e:
        print(f"Error generando terminos_condiciones.pdf: {e}")
        
    try:
        generar_catalogo_productos()
    except Exception as e:
        print(f"Error generando catalogo_productos.xlsx: {e}")
        
    try:
        generar_precios_envio()
    except Exception as e:
        print(f"Error generando precios_envio.csv: {e}")
        
    try:
        generar_config_tienda()
    except Exception as e:
        print(f"Error generando config_tienda.json: {e}")
        
    try:
        generar_presentacion_empresa()
    except Exception as e:
        print(f"Error generando presentacion_empresa.pptx: {e}")
        
    print("¡Generación de documentos finalizada!")

if __name__ == '__main__':
    generate_all()
